"""Tests for the export step.

These tests run varied Markdown through all three exporters and confirm that the
common parts (headings, lists, nested lists, tables, blockquotes, and code) are
not silently dropped. Export is where a hidden bug costs the user real rework,
so these are the most important tests in the project.

Run them with:

    pip install -e ".[dev]"
    pytest
"""

from __future__ import annotations

from pathlib import Path

import pytest

pytest.importorskip("markdown")

from markdown_mentor.export import (
    _parse_markdown,
    export_file,
    load_style,
)


VARIED = """# Title

Intro paragraph with **bold**, *italic*, and `code`.

## Table section

| Term | Meaning |
| --- | --- |
| Thesis | Main argument |
| Evidence | Support for a claim |

## Blockquote section

> A quote that runs
> across two lines.

## Nested list section

- Top item
  - Sub item one
  - Sub item two
- Second top item

1. First
2. Second
   1. Nested number

## Code section

```python
print("hello")
```
"""


def test_parser_finds_a_table():
    blocks = _parse_markdown(VARIED)
    tables = [b for b in blocks if b.kind == "table"]
    assert len(tables) == 1
    rows = tables[0].rows
    assert rows[0] == ["Term", "Meaning"]
    assert ["Thesis", "Main argument"] in rows
    assert ["Evidence", "Support for a claim"] in rows


def test_parser_finds_a_blockquote():
    blocks = _parse_markdown(VARIED)
    quotes = [b for b in blocks if b.kind == "quote"]
    assert len(quotes) == 1
    assert "across two lines" in quotes[0].text


def test_parser_records_nested_list_depth():
    blocks = _parse_markdown(VARIED)
    bullets = [b for b in blocks if b.kind == "bullet"]
    levels = {b.text: b.level for b in bullets}
    assert levels["Top item"] == 0
    assert levels["Sub item one"] == 1
    assert levels["Sub item two"] == 1
    nested_number = [b for b in blocks if b.kind == "number" and b.text == "Nested number"]
    assert nested_number and nested_number[0].level == 1


def test_parser_keeps_code_intact():
    blocks = _parse_markdown(VARIED)
    code = [b for b in blocks if b.kind == "code"]
    assert len(code) == 1
    assert 'print("hello")' in code[0].text


def test_table_pipes_do_not_leak_as_paragraphs():
    """The raw '| Term | Meaning |' text must not survive as a paragraph."""
    blocks = _parse_markdown(VARIED)
    paras = [b.text for b in blocks if b.kind == "para"]
    assert not any(p.strip().startswith("|") for p in paras)


def test_export_html_contains_table_and_quote(tmp_path):
    src = tmp_path / "m.md"
    src.write_text(VARIED, encoding="utf-8")
    out = export_file(src, "html", tmp_path / "m.html")
    html = Path(out).read_text(encoding="utf-8")
    assert "<table>" in html
    assert "<blockquote>" in html
    assert "Thesis" in html and "Main argument" in html


def test_export_docx_is_valid_and_has_a_table(tmp_path):
    from docx import Document

    src = tmp_path / "m.md"
    src.write_text(VARIED, encoding="utf-8")
    out = export_file(src, "docx", tmp_path / "m.docx")
    doc = Document(str(out))
    assert len(doc.tables) == 1
    cell_texts = [c.text for row in doc.tables[0].rows for c in row.cells]
    assert "Thesis" in cell_texts
    all_text = "\n".join(p.text for p in doc.paragraphs)
    assert "Title" in all_text
    assert "across two lines" in all_text  # blockquote survived


def test_export_pptx_is_valid_and_splits_slides(tmp_path):
    from pptx import Presentation

    src = tmp_path / "m.md"
    src.write_text(VARIED, encoding="utf-8")
    out = export_file(src, "pptx", tmp_path / "m.pptx")
    prs = Presentation(str(out))
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title]
    assert "Title" in titles
    assert "Table section" in titles
    # The table content should appear somewhere in a slide body.
    all_body = []
    for s in prs.slides:
        for shape in s.shapes:
            if shape.has_text_frame:
                all_body.append(shape.text_frame.text)
    joined = "\n".join(all_body)
    assert "Thesis" in joined


def test_pdf_export_works_or_refuses_clearly(tmp_path):
    """PDF either succeeds (if LibreOffice is present) or raises a clear,
    friendly error pointing to the 'Save as PDF' alternative."""
    from markdown_mentor.export import (
        LibreOfficeMissing,
        _find_libreoffice,
        export_file,
    )

    src = tmp_path / "m.md"
    src.write_text(VARIED, encoding="utf-8")

    if _find_libreoffice():
        out = export_file(src, "pdf", tmp_path / "m.pdf")
        data = Path(out).read_bytes()
        assert data[:4] == b"%PDF"
    else:
        with pytest.raises(LibreOfficeMissing) as info:
            export_file(src, "pdf", tmp_path / "m.pdf")
        assert "Save as PDF" in str(info.value)


def test_unknown_format_is_refused(tmp_path):
    src = tmp_path / "m.md"
    src.write_text("# Hi", encoding="utf-8")
    with pytest.raises(ValueError):
        export_file(src, "rtf", tmp_path / "m.rtf")


def test_markdown_style_file_changes_heading_size(tmp_path):
    from docx import Document

    style_file = tmp_path / "style.md"
    style_file.write_text("# Style\n\n- h1_size_pt: 30\n", encoding="utf-8")
    src = tmp_path / "m.md"
    src.write_text("# Heading\n\nText.", encoding="utf-8")
    out = export_file(src, "docx", tmp_path / "m.docx", style_file)
    doc = Document(str(out))
    for p in doc.paragraphs:
        if p.text == "Heading" and p.runs:
            # 30pt in EMU is 30 * 12700 = 381000
            assert p.runs[0].font.size == 381000
            break
    else:
        pytest.fail("Heading not found in document")


def test_partial_markdown_style_file_fills_defaults(tmp_path):
    style = load_style(None)
    assert style["body_font"]  # default present
    partial = tmp_path / "style.md"
    partial.write_text("# Style\n\n- body_font: Times New Roman\n", encoding="utf-8")
    merged = load_style(partial)
    assert merged["body_font"] == "Times New Roman"
    assert merged["h1_size_pt"] == style["h1_size_pt"]  # default kept
