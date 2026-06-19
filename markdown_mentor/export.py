"""Turn finished Markdown teaching materials into Word, PowerPoint, and web files.

Plain English: once you and the AI have produced teaching materials as Markdown,
this code converts each Markdown file into a format people actually use:

- DOCX  -> a Word document (good for handouts, worksheets, teacher notes)
- PPTX  -> a PowerPoint slide deck (good for lesson slides)
- HTML  -> a web page (good for sharing online)
- PDF   -> a PDF document, if LibreOffice is installed

Who does what:
- The user (with the AI) writes the teaching materials as Markdown.
- Markdown Mentor (this code) converts them to the chosen format.

A "style profile" controls how Markdown parts become document styles, such as
which heading sizes to use. Style profiles are plain text (JSON) files. A
default profile is built in, so you do not need one to get started.

PDF export uses LibreOffice. If LibreOffice is not installed, export to Word or
HTML and use 'Save as PDF' or 'Print to PDF' in your own software.
"""

from __future__ import annotations

import json
import re
from dataclasses import dataclass, field
from pathlib import Path

import markdown as md_lib
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt


DEFAULT_STYLE = {
    "name": "default",
    "body_font": "Calibri",
    "body_size_pt": 11,
    "heading_font": "Calibri",
    "heading_color": "1F3864",
    "h1_size_pt": 20,
    "h2_size_pt": 16,
    "h3_size_pt": 13,
    "html_title": "Teaching material",
}


@dataclass
class Block:
    """A single piece of a Markdown document, already classified.

    For most kinds, `text` holds the content and `level` holds a heading level
    or a list indent depth. For a "table" block, `rows` holds the parsed cells
    (the first row is the header).
    """

    kind: str   # "h1".."h3", "para", "bullet", "number", "code", "quote",
                # "table", "blank"
    text: str = ""
    level: int = 0
    rows: list[list[str]] = field(default_factory=list)


@dataclass
class ExportResult:
    inputs: list[Path] = field(default_factory=list)
    outputs: list[Path] = field(default_factory=list)


def load_style(style_path: str | Path | None) -> dict:
    """Load a style profile, falling back to the built-in default.

    Any keys missing from the file are filled in from the default, so a
    partial style profile is fine.
    """
    style = dict(DEFAULT_STYLE)
    if style_path:
        data = json.loads(Path(style_path).read_text(encoding="utf-8"))
        style.update(data)
    return style


def _looks_like_table_row(line: str) -> bool:
    """A table row has a pipe and is not a heading."""
    stripped = line.strip()
    return stripped.startswith("|") or ("|" in stripped and not stripped.startswith("#"))


def _is_table_divider(line: str) -> bool:
    """The line under a table header, made of dashes, colons, and pipes."""
    cells = [c.strip() for c in line.strip().strip("|").split("|")]
    return bool(cells) and all(re.fullmatch(r":?-{1,}:?", c) for c in cells if c)


def _split_row(line: str) -> list[str]:
    """Split one table row into its cell texts."""
    return [c.strip() for c in line.strip().strip("|").split("|")]


def _indent_depth(raw: str) -> int:
    """How deeply a list item is indented, in steps of two spaces (tabs = 2)."""
    leading = len(raw) - len(raw.lstrip(" \t"))
    spaces = raw[:leading].replace("\t", "  ")
    return len(spaces) // 2


def _parse_markdown(text: str) -> list[Block]:
    """Turn Markdown text into a simple list of blocks we can render.

    This is a deliberately small parser aimed at the Markdown that teaching
    materials normally use: headings, paragraphs, bullet and numbered lists
    (including nested ones), tables, blockquotes, and fenced code blocks. It
    does not try to handle every Markdown feature, but it should not silently
    drop the common ones.
    """
    blocks: list[Block] = []
    in_code = False
    code_lines: list[str] = []
    lines = text.splitlines()
    i = 0

    while i < len(lines):
        raw = lines[i].rstrip("\n")

        # Fenced code blocks come first: nothing inside them is parsed.
        if raw.strip().startswith("```"):
            if in_code:
                blocks.append(Block("code", "\n".join(code_lines)))
                code_lines = []
                in_code = False
            else:
                in_code = True
            i += 1
            continue
        if in_code:
            code_lines.append(raw)
            i += 1
            continue

        if not raw.strip():
            blocks.append(Block("blank"))
            i += 1
            continue

        heading = re.match(r"^(#{1,3})\s+(.*)$", raw)
        if heading:
            level = len(heading.group(1))
            blocks.append(Block(f"h{level}", heading.group(2).strip(), level))
            i += 1
            continue

        # Table: a row, then a divider row, then more rows.
        if (
            _looks_like_table_row(raw)
            and i + 1 < len(lines)
            and _is_table_divider(lines[i + 1])
        ):
            rows = [_split_row(raw)]
            i += 2  # skip header and divider
            while i < len(lines) and _looks_like_table_row(lines[i]) and lines[i].strip():
                rows.append(_split_row(lines[i]))
                i += 1
            blocks.append(Block("table", rows=rows))
            continue

        # Blockquote: one or more lines beginning with ">".
        quote = re.match(r"^\s*>\s?(.*)$", raw)
        if quote:
            quote_lines = [quote.group(1)]
            i += 1
            while i < len(lines):
                m = re.match(r"^\s*>\s?(.*)$", lines[i])
                if not m:
                    break
                quote_lines.append(m.group(1))
                i += 1
            blocks.append(Block("quote", " ".join(l.strip() for l in quote_lines).strip()))
            continue

        bullet = re.match(r"^(\s*)[-*+]\s+(.*)$", raw)
        if bullet:
            blocks.append(Block("bullet", bullet.group(2).strip(), _indent_depth(raw)))
            i += 1
            continue

        numbered = re.match(r"^(\s*)\d+[.)]\s+(.*)$", raw)
        if numbered:
            blocks.append(Block("number", numbered.group(2).strip(), _indent_depth(raw)))
            i += 1
            continue

        blocks.append(Block("para", raw.strip()))
        i += 1

    if in_code and code_lines:
        blocks.append(Block("code", "\n".join(code_lines)))
    return blocks


def _strip_inline(text: str) -> str:
    """Remove the most common inline Markdown marks for plain renderers."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"(?<!\*)\*(?!\*)(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\((.+?)\)", r"\1 (\2)", text)
    return text


# --- DOCX -----------------------------------------------------------------

def _to_docx(blocks: list[Block], out_path: Path, style: dict) -> None:
    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = style["body_font"]
    normal.font.size = Pt(style["body_size_pt"])

    color = RGBColor.from_string(style["heading_color"])
    heading_sizes = {1: style["h1_size_pt"], 2: style["h2_size_pt"], 3: style["h3_size_pt"]}

    for block in blocks:
        if block.kind.startswith("h"):
            p = doc.add_heading(level=block.level)
            run = p.add_run(_strip_inline(block.text))
            run.font.name = style["heading_font"]
            run.font.size = Pt(heading_sizes[block.level])
            run.font.color.rgb = color
        elif block.kind == "para":
            doc.add_paragraph(_strip_inline(block.text))
        elif block.kind == "bullet":
            # Word has List Bullet, List Bullet 2, List Bullet 3 for nesting.
            suffix = "" if block.level == 0 else f" {min(block.level + 1, 3)}"
            doc.add_paragraph(_strip_inline(block.text), style=f"List Bullet{suffix}")
        elif block.kind == "number":
            suffix = "" if block.level == 0 else f" {min(block.level + 1, 3)}"
            doc.add_paragraph(_strip_inline(block.text), style=f"List Number{suffix}")
        elif block.kind == "quote":
            p = doc.add_paragraph(_strip_inline(block.text), style="Intense Quote")
        elif block.kind == "table":
            _docx_table(doc, block.rows, style)
        elif block.kind == "code":
            p = doc.add_paragraph()
            run = p.add_run(block.text)
            run.font.name = "Consolas"
            run.font.size = Pt(style["body_size_pt"] - 1)
        # blank lines need no action in Word

    doc.save(str(out_path))


def _docx_table(doc, rows: list[list[str]], style: dict) -> None:
    """Add a table to the Word document. The first row is the header."""
    if not rows:
        return
    cols = max(len(r) for r in rows)
    table = doc.add_table(rows=0, cols=cols)
    table.style = "Light Grid Accent 1"
    for r_index, row in enumerate(rows):
        cells = table.add_row().cells
        for c_index in range(cols):
            text = row[c_index] if c_index < len(row) else ""
            cell = cells[c_index]
            cell.text = ""
            run = cell.paragraphs[0].add_run(_strip_inline(text))
            run.font.size = Pt(style["body_size_pt"])
            if r_index == 0:
                run.bold = True


# --- PPTX -----------------------------------------------------------------

def _to_pptx(blocks: list[Block], out_path: Path, style: dict) -> None:
    """Build a slide deck. Each H1 or H2 starts a new slide; the text and
    bullets under it become that slide's body. Tables and quotes are flattened
    into readable lines so nothing is dropped."""
    prs = Presentation()
    blank_title_content = prs.slide_layouts[1]

    # Each body entry is (text, indent_level).
    slides: list[tuple[str, list[tuple[str, int]]]] = []
    current_title = None
    current_body: list[tuple[str, int]] = []

    def flush():
        if current_title is not None or current_body:
            slides.append((current_title or "Slide", list(current_body)))

    for block in blocks:
        if block.kind in ("h1", "h2"):
            flush()
            current_title = _strip_inline(block.text)
            current_body = []
        elif block.kind == "h3":
            current_body.append((_strip_inline(block.text), 0))
        elif block.kind in ("para", "number", "bullet"):
            current_body.append((_strip_inline(block.text), block.level))
        elif block.kind == "quote":
            current_body.append(("\u201c" + _strip_inline(block.text) + "\u201d", 0))
        elif block.kind == "table":
            for row in block.rows:
                current_body.append((" | ".join(_strip_inline(c) for c in row), 0))
        elif block.kind == "code":
            for code_line in block.text.splitlines():
                current_body.append((code_line, 0))
    flush()

    if not slides:
        slides = [("Slide", [])]

    for title, body in slides:
        slide = prs.slides.add_slide(blank_title_content)
        slide.shapes.title.text = title
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        if body:
            first_text, first_level = body[0]
            tf.paragraphs[0].text = first_text
            tf.paragraphs[0].level = min(first_level, 4)
            for text, level in body[1:]:
                p = tf.add_paragraph()
                p.text = text
                p.level = min(level, 4)
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = PptPt(18)

    prs.save(str(out_path))


# --- HTML -----------------------------------------------------------------

def _to_html(markdown_text: str, out_path: Path, style: dict) -> None:
    body = md_lib.markdown(
        markdown_text,
        extensions=["fenced_code", "tables", "sane_lists"],
    )
    title = style.get("html_title", "Teaching material")
    page = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{title}</title>
<style>
  body {{ font-family: {style['body_font']}, system-ui, sans-serif;
         max-width: 50rem; margin: 2rem auto; padding: 0 1rem;
         line-height: 1.6; color: #1a1a1a; }}
  h1, h2, h3 {{ font-family: {style['heading_font']}, system-ui, sans-serif;
               color: #{style['heading_color']}; line-height: 1.25; }}
  code {{ background: #f2f2f2; padding: 0.1em 0.3em; border-radius: 3px; }}
  pre {{ background: #f2f2f2; padding: 1rem; overflow-x: auto; border-radius: 6px; }}
  table {{ border-collapse: collapse; }}
  th, td {{ border: 1px solid #ccc; padding: 0.4rem 0.6rem; }}
</style>
</head>
<body>
{body}
</body>
</html>
"""
    out_path.write_text(page, encoding="utf-8")


# --- public entry point ---------------------------------------------------

VALID_FORMATS = {"docx", "pptx", "html", "pdf"}


class LibreOfficeMissing(RuntimeError):
    """Raised when PDF export is asked for but LibreOffice is not installed."""


def _find_libreoffice() -> str | None:
    """Return the LibreOffice command if it is installed, else None.

    PDF export uses LibreOffice to convert a Word document into a PDF. It is an
    optional extra: most users can use 'Save as PDF' in their own software
    instead, so LibreOffice is not a required dependency.
    """
    import shutil

    for name in ("libreoffice", "soffice"):
        found = shutil.which(name)
        if found:
            return found
    return None


def _docx_to_pdf(docx_path: Path, output_path: Path) -> None:
    """Convert a Word document to PDF using LibreOffice headless."""
    import subprocess
    import tempfile

    soffice = _find_libreoffice()
    if not soffice:
        raise LibreOfficeMissing(
            "PDF export needs LibreOffice, and it was not found on your "
            "computer.\n\n"
            "You have two choices:\n"
            "  1. Export to Word (docx) or web (html), open the file, and use "
            "'Save as PDF' in your own software. This needs nothing extra.\n"
            "  2. Install LibreOffice (free) from https://www.libreoffice.org "
            "and try PDF export again."
        )

    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(docx_path)],
            check=True,
            capture_output=True,
            timeout=120,
        )
        produced = Path(tmp) / (docx_path.stem + ".pdf")
        if not produced.is_file():
            raise RuntimeError("LibreOffice ran but did not produce a PDF.")
        output_path.write_bytes(produced.read_bytes())


def export_file(
    input_path: str | Path,
    fmt: str,
    output_path: str | Path | None = None,
    style_path: str | Path | None = None,
) -> Path:
    """Convert one Markdown file to one output format.

    Arguments:
        input_path: the Markdown (.md) file to convert.
        fmt: one of "docx", "pptx", "html", or "pdf".
        output_path: where to write the result. Defaults to the same name with
            the new extension.
        style_path: an optional style profile (JSON) file.

    Returns the path of the file that was written.

    PDF export uses LibreOffice if it is installed. If it is not, a clear
    message explains the 'Save as PDF' alternative.
    """
    fmt = fmt.lower().lstrip(".")
    if fmt not in VALID_FORMATS:
        raise ValueError(
            f"Unknown format '{fmt}'. Choose one of: docx, pptx, html, pdf."
        )

    input_path = Path(input_path).expanduser().resolve()
    if not input_path.is_file():
        raise FileNotFoundError(f"Markdown file not found: {input_path}")

    text = input_path.read_text(encoding="utf-8")
    style = load_style(style_path)

    if output_path is None:
        output_path = input_path.with_suffix("." + fmt)
    output_path = Path(output_path).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if fmt == "html":
        _to_html(text, output_path, style)
    elif fmt == "pdf":
        # Render to a temporary Word document, then convert that to PDF.
        import tempfile

        blocks = _parse_markdown(text)
        with tempfile.TemporaryDirectory() as tmp:
            tmp_docx = Path(tmp) / (input_path.stem + ".docx")
            _to_docx(blocks, tmp_docx, style)
            _docx_to_pdf(tmp_docx, output_path)
    else:
        blocks = _parse_markdown(text)
        if fmt == "docx":
            _to_docx(blocks, output_path, style)
        else:
            _to_pptx(blocks, output_path, style)

    return output_path


def export_path(
    input_path: str | Path,
    fmt: str,
    output_dir: str | Path | None = None,
    style_path: str | Path | None = None,
) -> ExportResult:
    """Convert a single Markdown file, or every .md file in a folder.

    This is what the command line calls. It returns a record of inputs and
    outputs so the command can report clearly.
    """
    input_path = Path(input_path).expanduser().resolve()
    result = ExportResult()

    if input_path.is_dir():
        md_files = sorted(input_path.glob("*.md"))
    else:
        md_files = [input_path]

    for md_file in md_files:
        if output_dir:
            out = Path(output_dir).expanduser().resolve() / (md_file.stem + "." + fmt.lower())
        else:
            out = None
        written = export_file(md_file, fmt, out, style_path)
        result.inputs.append(md_file)
        result.outputs.append(written)

    return result
