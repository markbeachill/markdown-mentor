#!/usr/bin/env python3
"""Guide and export Markdown Mentor teaching materials.

Normal use from inside a Markdown Mentor project folder:

    python make-teaching-materials.py guide
    python make-teaching-materials.py export html
    python make-teaching-materials.py export docx

This script does not send files to AI. It shows the workflow and converts the
Markdown files that you save in 5-draft-materials/.
"""

from __future__ import annotations

import argparse
import html
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

DEFAULT_INPUT = Path("5-draft-materials")
DEFAULT_OUTPUT = Path("6-final-exports")
DEFAULT_STYLE = Path("style/style.md")


def load_style(path: Path | None) -> dict[str, str]:
    style = {
        "body_font": "Calibri",
        "body_size_pt": "11",
        "heading_font": "Calibri",
        "heading_color": "1F3864",
        "h1_size_pt": "20",
        "h2_size_pt": "16",
        "h3_size_pt": "13",
        "title_size_pt": "32",
        "slide_body_size_pt": "20",
        "html_title": "Teaching material",
    }
    if path and path.is_file():
        for line in path.read_text(encoding="utf-8", errors="replace").splitlines():
            m = re.match(r"\s*-\s*([A-Za-z0-9_ -]+)\s*:\s*(.+?)\s*$", line)
            if m:
                key = m.group(1).strip().replace(" ", "_")
                style[key] = m.group(2).strip()
    return style


def markdown_files(path: Path) -> list[Path]:
    if path.is_file():
        return [path] if path.suffix.lower() in {".md", ".markdown"} else []
    if path.is_dir():
        return sorted([p for p in path.rglob("*") if p.suffix.lower() in {".md", ".markdown"}])
    raise FileNotFoundError(f"Not found: {path}")


def split_blocks(text: str) -> list[tuple[str, str]]:
    blocks: list[tuple[str, str]] = []
    para: list[str] = []

    def flush() -> None:
        nonlocal para
        if para:
            blocks.append(("p", " ".join(x.strip() for x in para).strip()))
            para = []

    in_code = False
    code: list[str] = []
    for raw in text.splitlines():
        line = raw.rstrip("\n")
        if line.strip().startswith("```"):
            if in_code:
                blocks.append(("code", "\n".join(code)))
                code = []
                in_code = False
            else:
                flush()
                in_code = True
            continue
        if in_code:
            code.append(line)
            continue
        if not line.strip():
            flush()
            continue
        m = re.match(r"^(#{1,6})\s+(.+)$", line)
        if m:
            flush()
            blocks.append((f"h{len(m.group(1))}", m.group(2).strip()))
        elif re.match(r"^\s*[-*+]\s+", line):
            flush()
            blocks.append(("li", re.sub(r"^\s*[-*+]\s+", "", line).strip()))
        elif re.match(r"^\s*\d+[.)]\s+", line):
            flush()
            blocks.append(("li", re.sub(r"^\s*\d+[.)]\s+", "", line).strip()))
        elif line.startswith(">"):
            flush()
            blocks.append(("quote", line.lstrip("> ").strip()))
        else:
            para.append(line)
    flush()
    return blocks


def strip_inline(text: str) -> str:
    text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1", text)
    text = re.sub(r"[*_`~]", "", text)
    return text


def to_html(markdown_text: str, out_path: Path, style: dict[str, str]) -> None:
    css = f"""
body {{ font-family: {style.get('body_font', 'Calibri')}, sans-serif; line-height: 1.55; max-width: 860px; margin: 2rem auto; padding: 0 1rem; }}
h1, h2, h3 {{ color: #{style.get('heading_color', '1F3864')}; }}
pre {{ background: #f6f8fa; padding: 1rem; overflow-x: auto; }}
blockquote {{ border-left: 4px solid #ddd; padding-left: 1rem; color: #555; }}
"""
    parts = ["<!doctype html>", "<html lang='en'>", "<head>", "<meta charset='utf-8'>", f"<title>{html.escape(style.get('html_title','Teaching material'))}</title>", f"<style>{css}</style>", "</head><body>"]
    in_ul = False
    for kind, value in split_blocks(markdown_text):
        safe = html.escape(strip_inline(value))
        if kind != "li" and in_ul:
            parts.append("</ul>")
            in_ul = False
        if kind.startswith("h") and kind[1:].isdigit():
            parts.append(f"<{kind}>{safe}</{kind}>")
        elif kind == "li":
            if not in_ul:
                parts.append("<ul>")
                in_ul = True
            parts.append(f"<li>{safe}</li>")
        elif kind == "quote":
            parts.append(f"<blockquote>{safe}</blockquote>")
        elif kind == "code":
            parts.append(f"<pre><code>{html.escape(value)}</code></pre>")
        else:
            parts.append(f"<p>{safe}</p>")
    if in_ul:
        parts.append("</ul>")
    parts.append("</body></html>")
    out_path.write_text("\n".join(parts), encoding="utf-8")


def to_docx(markdown_text: str, out_path: Path, style: dict[str, str]) -> None:
    try:
        from docx import Document  # type: ignore
        from docx.shared import Pt, RGBColor  # type: ignore
    except ImportError as exc:
        raise RuntimeError("DOCX export needs python-docx. Install it with: python -m pip install python-docx") from exc
    doc = Document()
    for kind, value in split_blocks(markdown_text):
        text = strip_inline(value)
        if kind == "h1":
            p = doc.add_heading(text, level=1)
        elif kind == "h2":
            p = doc.add_heading(text, level=2)
        elif kind == "h3":
            p = doc.add_heading(text, level=3)
        elif kind == "li":
            p = doc.add_paragraph(text, style="List Bullet")
        elif kind == "quote":
            p = doc.add_paragraph(text)
            p.paragraph_format.left_indent = Pt(18)
        elif kind == "code":
            p = doc.add_paragraph(value)
        else:
            p = doc.add_paragraph(text)
        for run in p.runs:
            if kind.startswith("h"):
                run.font.name = style.get("heading_font", "Calibri")
                try:
                    run.font.color.rgb = RGBColor.from_string(style.get("heading_color", "1F3864"))
                except ValueError:
                    pass
            else:
                run.font.name = style.get("body_font", "Calibri")
                try:
                    run.font.size = Pt(float(style.get("body_size_pt", "11")))
                except ValueError:
                    pass
    doc.save(str(out_path))


def to_pptx(markdown_text: str, out_path: Path, style: dict[str, str]) -> None:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Pt  # type: ignore
    except ImportError as exc:
        raise RuntimeError("PPTX export needs python-pptx. Install it with: python -m pip install python-pptx") from exc
    prs = Presentation()
    title = "Teaching material"
    chunks: list[tuple[str, list[str]]] = []
    current_title = None
    current_body: list[str] = []
    for kind, value in split_blocks(markdown_text):
        if kind in {"h1", "h2"}:
            if current_title:
                chunks.append((current_title, current_body))
            current_title = strip_inline(value)
            current_body = []
            if title == "Teaching material":
                title = current_title
        else:
            current_body.append(strip_inline(value))
    if current_title:
        chunks.append((current_title, current_body))
    if not chunks:
        chunks = [(title, [strip_inline(v) for _k, v in split_blocks(markdown_text)])]
    for slide_title, body in chunks:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(body[:8]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            try:
                p.font.size = Pt(float(style.get("slide_body_size_pt", style.get("body_size_pt", "20"))))
            except ValueError:
                pass
    prs.save(str(out_path))


def find_libreoffice() -> str | None:
    for name in ["soffice", "libreoffice"]:
        path = shutil.which(name)
        if path:
            return path
    return None


def docx_to_pdf(docx_path: Path, pdf_path: Path) -> None:
    office = find_libreoffice()
    if not office:
        raise RuntimeError("PDF export needs LibreOffice. Install LibreOffice or export DOCX and save as PDF manually.")
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run([office, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(docx_path)], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        made = Path(tmp) / (docx_path.stem + ".pdf")
        if not made.exists():
            raise RuntimeError("LibreOffice did not create the PDF file.")
        shutil.move(str(made), str(pdf_path))


def export_one(src: Path, fmt: str, out_dir: Path, style: dict[str, str]) -> Path:
    text = src.read_text(encoding="utf-8", errors="replace")
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / f"{src.stem}.{fmt}"
    if fmt == "html":
        to_html(text, out, style)
    elif fmt == "docx":
        to_docx(text, out, style)
    elif fmt == "pptx":
        to_pptx(text, out, style)
    elif fmt == "pdf":
        with tempfile.TemporaryDirectory() as tmp:
            tmp_docx = Path(tmp) / f"{src.stem}.docx"
            to_docx(text, tmp_docx, style)
            docx_to_pdf(tmp_docx, out)
    else:
        raise ValueError(f"Unknown format: {fmt}")
    return out


def command_guide(_args: argparse.Namespace) -> int:
    print("Markdown Mentor guide")
    print()
    steps = [
        "1. Put source files in 1-source-files/.",
        "2. Run: python make-markdown-library.py make",
        "3. Upload 2-markdown-library/markdown-library.md to AI.",
        "4. Use 3-teaching-approach/prompt-create-teaching-approach.md.",
        "5. Save the final approach as 3-teaching-approach/teaching-approach.md.",
        "6. Use 3-teaching-approach/prompt-check-teaching-approach.md.",
        "7. Edit 4-teaching-materials-pack/teaching-materials-pack.md.",
        "8. Use 4-teaching-materials-pack/prompt-create-teaching-materials.md.",
        "9. Save AI-created Markdown files in 5-draft-materials/.",
        "10. Edit style/style.md if needed.",
        "11. Run: python make-teaching-materials.py export docx",
    ]
    for step in steps:
        print(step)
    return 0


def command_export(args: argparse.Namespace) -> int:
    fmt = args.format.lower()
    if fmt not in {"docx", "pptx", "html", "pdf"}:
        print("Choose one format: docx, pptx, html, or pdf")
        return 1
    inp = Path(args.input or DEFAULT_INPUT)
    out = Path(args.output or DEFAULT_OUTPUT)
    style = load_style(Path(args.style or DEFAULT_STYLE))
    try:
        files = markdown_files(inp)
        if not files:
            print(f"No Markdown files found in {inp}.")
            return 1
        for src in files:
            dst = export_one(src, fmt, out, style)
            print(f"{src} -> {dst}")
    except (FileNotFoundError, RuntimeError, ValueError, subprocess.CalledProcessError) as exc:
        print(f"Problem: {exc}")
        return 1
    return 0


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Guide and export Markdown Mentor teaching materials.")
    sub = parser.add_subparsers(dest="command", required=True)
    p = sub.add_parser("guide", help="Show the project workflow steps.")
    p.set_defaults(func=command_guide)
    p = sub.add_parser("export", help="Export Markdown teaching materials.")
    p.add_argument("format", help="docx, pptx, html, or pdf")
    p.add_argument("--input", help="Input Markdown file or folder. Default: 5-draft-materials")
    p.add_argument("--output", help="Output folder. Default: 6-final-exports")
    p.add_argument("--style", help="Style file. Default: style/style.md")
    p.set_defaults(func=command_export)
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    return args.func(args)


if __name__ == "__main__":
    raise SystemExit(main())
